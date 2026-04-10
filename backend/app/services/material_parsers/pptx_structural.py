"""
Structural PPTX parser — converts a PowerPoint file into TipTap content_json.

This is the v1 default parser for the ``.pptx`` format. It walks slides via
python-pptx and produces:

- A ``heading`` (level 1) per slide title
- ``paragraph``, ``bulletList``, and ``table`` nodes from text/table shapes
- ``image`` nodes (with placeholder ``src`` filenames) plus
  ``ExtractedImage`` entries for each picture shape
- A ``speakerNotes`` node per slide that has notes
- ``slideBreak`` nodes between slides

Per ADR-038 the parser deliberately discards anything not in the slide's
content layer: animations, transitions, slide masters, decorative shapes,
backgrounds. Per ADR-061 the parser emits warnings rather than silently
dropping content the user might miss (deeply nested bullets, equations,
linked images that can't be embedded, etc.).

Per ADR-064 speaker notes ride along here as first-class structured nodes
(not flattened to text), then the export pipeline routes them through
Pandoc's ``::: notes`` fenced div on PPTX export to round-trip back to
PowerPoint's speaker pane.
"""

from __future__ import annotations

import io
import logging
from pathlib import Path
from typing import Any, ClassVar

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.services.material_parsers.base import (
    ExtractedImage,
    MaterialParser,
    MaterialParseResult,
)

logger = logging.getLogger(__name__)


# Maximum bullet nesting depth we emit. Beyond this we flatten and warn.
# TipTap nested bulletLists are valid but visually noisy beyond 1-2 levels;
# educators rarely use deep nesting in real lecture content, so flattening
# is a safe v1 simplification.
_MAX_BULLET_LEVEL = 1


class PptxStructuralParser(MaterialParser):
    """Convert a PowerPoint deck into TipTap content_json with slide breaks."""

    name: ClassVar[str] = "pptx_structural"
    display_name: ClassVar[str] = "PowerPoint (structural)"
    description: ClassVar[str] = (
        "Default PPTX importer. Extracts slide titles, body text, "
        "bullet lists, tables, images, and speaker notes into editable "
        "structured content. Discards animations, slide masters, and "
        "decorative shapes."
    )
    supported_formats: ClassVar[list[str]] = ["pptx"]
    requires_ai: ClassVar[bool] = False

    async def parse(
        self,
        file_content: bytes,
        filename: str,
        *,
        user_context: dict[str, object] | None = None,
    ) -> MaterialParseResult:
        # user_context reserved for sector-aware parsing in later phases
        _ = user_context

        prs = Presentation(io.BytesIO(file_content))

        top_nodes: list[dict[str, Any]] = []
        images: list[ExtractedImage] = []
        warnings: list[str] = []
        deck_title: str | None = None
        skipped_slides = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_nodes, slide_images, slide_warnings = self._parse_slide(
                slide, slide_num
            )

            # Capture the first non-empty slide title as the deck title
            if deck_title is None:
                first_heading = next(
                    (n for n in slide_nodes if n.get("type") == "heading"), None
                )
                if first_heading:
                    deck_title = self._node_text(first_heading)

            # Drop slides that produced absolutely nothing (no title, no body,
            # no notes, no images). Rare in practice but possible.
            if not slide_nodes:
                skipped_slides += 1
                continue

            # Slide breaks separate slides — insert before every slide
            # except the first emitted one (so leading slideBreak is avoided)
            if top_nodes:
                top_nodes.append({"type": "slideBreak"})

            top_nodes.extend(slide_nodes)
            images.extend(slide_images)
            warnings.extend(slide_warnings)

        if skipped_slides:
            warnings.append(
                f"{skipped_slides} empty slide(s) skipped (no title, body, or notes)"
            )

        return MaterialParseResult(
            title=deck_title or Path(filename).stem,
            content_json={"type": "doc", "content": top_nodes},
            images=images,
            warnings=warnings,
            parser_used=self.name,
        )

    # ------------------------------------------------------------------
    # Slide-level parsing
    # ------------------------------------------------------------------

    def _parse_slide(
        self, slide: Any, slide_num: int
    ) -> tuple[list[dict[str, Any]], list[ExtractedImage], list[str]]:
        nodes: list[dict[str, Any]] = []
        images: list[ExtractedImage] = []
        warnings: list[str] = []

        # Title placeholder. python-pptx exposes this via slide.shapes.title
        # which may be None for slides with no title placeholder.
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except (AttributeError, KeyError):
            title_shape = None

        title_shape_id = title_shape.shape_id if title_shape is not None else None

        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()
            if title_text:
                nodes.append(
                    {
                        "type": "heading",
                        "attrs": {"level": 1},
                        "content": [{"type": "text", "text": title_text}],
                    }
                )

        # Walk all other shapes in document order
        img_counter = [0]  # mutable counter so helpers can increment
        self._walk_shapes(
            slide.shapes,
            slide_num=slide_num,
            skip_shape_id=title_shape_id,
            nodes=nodes,
            images=images,
            warnings=warnings,
            img_counter=img_counter,
        )

        # Speaker notes — per ADR-064 these are first-class structured
        # content, not flattened to body text
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                nodes.append(self._speaker_notes_node(notes_text))

        return nodes, images, warnings

    def _walk_shapes(
        self,
        shapes: Any,
        *,
        slide_num: int,
        skip_shape_id: int | None,
        nodes: list[dict[str, Any]],
        images: list[ExtractedImage],
        warnings: list[str],
        img_counter: list[int],
    ) -> None:
        """Walk a shape collection (slide shapes or grouped child shapes).

        Recursive: ``MSO_SHAPE_TYPE.GROUP`` shapes are flattened by
        descending into their children.
        """
        for shape in shapes:
            if skip_shape_id is not None and shape.shape_id == skip_shape_id:
                continue

            shape_type = shape.shape_type

            if shape_type == MSO_SHAPE_TYPE.GROUP:
                # Recurse into grouped shapes — group is just a container
                self._walk_shapes(
                    shape.shapes,
                    slide_num=slide_num,
                    skip_shape_id=None,
                    nodes=nodes,
                    images=images,
                    warnings=warnings,
                    img_counter=img_counter,
                )
                continue

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                self._handle_picture(
                    shape, slide_num, img_counter, nodes, images, warnings
                )
                continue

            if shape.has_text_frame:
                nodes.extend(self._text_frame_to_nodes(shape.text_frame, warnings))
                continue

            if shape.has_table:
                nodes.append(self._table_to_node(shape.table))
                continue

            # Anything else (charts, SmartArt without text, OLE objects) is
            # silently skipped per ADR-038 — these are visual decoration
            # rather than content we can meaningfully extract.

    # ------------------------------------------------------------------
    # Picture shapes → image nodes + ExtractedImage
    # ------------------------------------------------------------------

    def _handle_picture(
        self,
        shape: Any,
        slide_num: int,
        img_counter: list[int],
        nodes: list[dict[str, Any]],
        images: list[ExtractedImage],
        warnings: list[str],
    ) -> None:
        # Linked (non-embedded) pictures raise on .image access. Skip with
        # a warning rather than crashing the import.
        try:
            image_obj = shape.image
        except (AttributeError, ValueError):
            warnings.append(
                f"Slide {slide_num}: linked image skipped (only embedded images are supported)"
            )
            return

        ext = (image_obj.ext or "png").lower()
        idx = img_counter[0]
        filename = f"slide-{slide_num}-{idx}.{ext}"
        img_counter[0] += 1

        images.append(
            ExtractedImage(
                filename=filename,
                data=image_obj.blob,
                mime_type=image_obj.content_type,
            )
        )
        nodes.append(
            {
                "type": "image",
                "attrs": {
                    "src": filename,  # Placeholder; apply route rewrites to canonical URL
                    "alt": "",
                },
            }
        )

    # ------------------------------------------------------------------
    # Text frames → paragraphs / bullet lists
    # ------------------------------------------------------------------

    def _text_frame_to_nodes(
        self, text_frame: Any, warnings: list[str]
    ) -> list[dict[str, Any]]:
        """Convert a python-pptx TextFrame to a sequence of TipTap nodes.

        Consecutive paragraphs at bullet level >= 1 are gathered into a
        single ``bulletList``. Bullet nesting beyond ``_MAX_BULLET_LEVEL``
        is flattened (with a warning) — TipTap supports nested bulletLists
        but real lecture content rarely needs more than one level.
        """
        nodes: list[dict[str, Any]] = []
        current_list: dict[str, Any] | None = None
        deep_nesting_warned = False

        for para in text_frame.paragraphs:
            text_nodes = self._runs_to_text_nodes(para)

            # Skip entirely empty paragraphs (no runs, no text)
            if not text_nodes:
                continue

            level = para.level or 0

            if level == 0:
                # Flush any open list before starting a plain paragraph
                if current_list is not None:
                    nodes.append(current_list)
                    current_list = None
                nodes.append({"type": "paragraph", "content": text_nodes})
            else:
                if level > _MAX_BULLET_LEVEL and not deep_nesting_warned:
                    warnings.append(
                        f"Bullet nesting beyond level {_MAX_BULLET_LEVEL} "
                        f"flattened (PPTX had level {level})"
                    )
                    deep_nesting_warned = True
                if current_list is None:
                    current_list = {"type": "bulletList", "content": []}
                current_list["content"].append(
                    {
                        "type": "listItem",
                        "content": [
                            {"type": "paragraph", "content": text_nodes}
                        ],
                    }
                )

        if current_list is not None:
            nodes.append(current_list)

        return nodes

    def _runs_to_text_nodes(self, paragraph: Any) -> list[dict[str, Any]]:
        """Convert a paragraph's runs into TipTap text nodes with marks.

        Empty runs are skipped. Tri-state font properties (True/False/None)
        are treated as "explicit True only" — None means "inherit" in
        python-pptx and we don't try to resolve theme defaults.
        """
        nodes: list[dict[str, Any]] = []

        for run in paragraph.runs:
            if not run.text:
                continue

            marks: list[dict[str, Any]] = []
            font = run.font
            if font.bold is True:
                marks.append({"type": "bold"})
            if font.italic is True:
                marks.append({"type": "italic"})
            if font.underline is True:
                marks.append({"type": "underline"})

            text_node: dict[str, Any] = {"type": "text", "text": run.text}
            if marks:
                text_node["marks"] = marks
            nodes.append(text_node)

        # Fallback: paragraph has text but no runs (rare). Emit a single
        # plain text node so we don't lose content.
        if not nodes and paragraph.text.strip():
            nodes.append({"type": "text", "text": paragraph.text})

        return nodes

    # ------------------------------------------------------------------
    # Tables
    # ------------------------------------------------------------------

    def _table_to_node(self, table: Any) -> dict[str, Any]:
        """Convert a PPTX table to a TipTap table node.

        Row 0 is treated as the header row. Cell text is extracted as plain
        text — formatting marks within cells are not preserved in v1
        (PPTX cell formatting is rarely meaningful for educational tables).
        """
        rows: list[dict[str, Any]] = []

        for row_idx, row in enumerate(table.rows):
            cells: list[dict[str, Any]] = []
            cell_type = "tableHeader" if row_idx == 0 else "tableCell"

            for cell in row.cells:
                cell_text = cell.text.strip()
                cell_content: list[dict[str, Any]] = []
                if cell_text:
                    cell_content = [
                        {
                            "type": "paragraph",
                            "content": [{"type": "text", "text": cell_text}],
                        }
                    ]
                else:
                    cell_content = [{"type": "paragraph", "content": []}]
                cells.append({"type": cell_type, "content": cell_content})

            rows.append({"type": "tableRow", "content": cells})

        return {"type": "table", "content": rows}

    # ------------------------------------------------------------------
    # Speaker notes
    # ------------------------------------------------------------------

    def _speaker_notes_node(self, text: str) -> dict[str, Any]:
        """Build a speakerNotes node from raw notes text.

        PPTX notes come back as a single string with newline-separated
        paragraphs. We split on newlines and emit one paragraph per line,
        skipping blank lines. Plain-text fidelity only — formatting in
        notes is unusual and not preserved in v1.
        """
        paragraphs = [
            line for line in (raw.strip() for raw in text.split("\n")) if line
        ]

        if not paragraphs:
            # Defensive: caller already checked text was non-empty
            return {
                "type": "speakerNotes",
                "content": [{"type": "paragraph", "content": []}],
            }

        return {
            "type": "speakerNotes",
            "content": [
                {
                    "type": "paragraph",
                    "content": [{"type": "text", "text": p}],
                }
                for p in paragraphs
            ],
        }

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _node_text(node: dict[str, Any]) -> str:
        """Extract plain text from a node tree (for title fallback)."""
        if node.get("type") == "text":
            return str(node.get("text", ""))
        return "".join(
            PptxStructuralParser._node_text(child)
            for child in node.get("content", [])
        )
