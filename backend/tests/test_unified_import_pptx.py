"""
Integration test for unified_import_service with structural PPTX parser.

Verifies that bulk LMS imports (zip files containing .pptx files) route
through the new pptx_structural parser and produce WeeklyMaterials with
populated content_json — not the legacy plain-text description that the
upstream code path produces for unstructured formats.

This is the consistency test the user asked for: regardless of whether
a PPTX comes in via the single-file Mode A endpoint or as part of a
bulk LMS / zip import, it should land in the database with the same
structured shape.
"""

from __future__ import annotations

import io
import zipfile
from typing import TYPE_CHECKING

import pytest
from pptx import Presentation

from app.models.weekly_material import WeeklyMaterial
from app.services import git_content_service
from app.services.unified_import_service import UnifiedImportService

if TYPE_CHECKING:
    from pathlib import Path

    from sqlalchemy.orm import Session

    from app.models.unit import Unit
    from app.models.user import User


# ---------------------------------------------------------------------------
# Builders
# ---------------------------------------------------------------------------


def _build_pptx(title: str, body_text: str, notes_text: str) -> bytes:
    """Construct a minimal one-slide PPTX with title, body, and notes."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = body_text
    slide.notes_slide.notes_text_frame.text = notes_text
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _build_zip_with_pptx(pptx_bytes: bytes, filename: str = "lecture.pptx") -> bytes:
    """Wrap a single PPTX in a plain zip (no manifest, so no LMS detection)."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr(filename, pptx_bytes)
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# Test
# ---------------------------------------------------------------------------


class TestUnifiedImportPptxStructured:
    """Bulk-import path must use the structural parser for .pptx files."""

    @pytest.mark.asyncio
    async def test_pptx_in_zip_produces_structured_content_json(
        self,
        test_db: Session,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        # Point the git singleton at a temp dir so image writes don't
        # pollute the real content_repos directory
        original_singleton = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )

        try:
            pptx = _build_pptx(
                title="Imported Lecture",
                body_text="First key point",
                notes_text="Pause here for questions",
            )
            zip_bytes = _build_zip_with_pptx(pptx, "lecture_01.pptx")

            svc = UnifiedImportService()
            result = svc.apply(
                zip_bytes,
                user_id=str(test_user.id),
                user_email=test_user.email,
                db=test_db,
                unit_code="TEST101",
                unit_title="Test Unit",
                duration_weeks=1,
            )

            # The apply method kicks off a background task; await it
            # explicitly so the test sees the persisted state
            await svc._bg_task

            # Find the unit that was created
            from app.models.unit import Unit

            unit = (
                test_db.query(Unit).filter(Unit.code == "TEST101").first()
            )
            assert unit is not None, "unified import should have created a unit"

            # And the material from the PPTX
            materials = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == unit.id)
                .all()
            )
            assert len(materials) == 1, (
                f"expected one material, got {len(materials)}"
            )
            mat = materials[0]

            # The structured parser path must populate content_json,
            # not just description (the regression we're guarding against)
            assert mat.content_json is not None, (
                "PPTX in bulk zip should produce structured content_json, "
                "not fall back to plain-text description"
            )
            assert mat.content_json["type"] == "doc"

            # Speaker notes survive the bulk-import path the same way
            # they do for the single-file Mode A path
            content_str = str(mat.content_json)
            assert "Pause here for questions" in content_str
            assert "First key point" in content_str
            assert "Imported Lecture" in content_str

            # And the speakerNotes node is structured (not flattened
            # into the body as "Notes: ..." text)
            assert "speakerNotes" in content_str

            _ = result  # task_id is informational only for this test
        finally:
            git_content_service._git_service = original_singleton

    @pytest.mark.asyncio
    async def test_pptx_in_zip_round_trips_to_pandoc_target(
        self,
        test_db: Session,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        """Closes the consistency loop end-to-end: a PPTX imported via the
        bulk LMS path produces content_json that the export pipeline can
        route back to PowerPoint's speaker notes pane via Pandoc fenced divs.
        Same assertion as the single-file Mode A round-trip test, but
        coming through the bulk import path."""
        original_singleton = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )

        try:
            pptx = _build_pptx(
                title="Round Trip Test",
                body_text="Body content",
                notes_text="Delivery prompt for the lecturer",
            )
            zip_bytes = _build_zip_with_pptx(pptx)

            svc = UnifiedImportService()
            svc.apply(
                zip_bytes,
                user_id=str(test_user.id),
                user_email=test_user.email,
                db=test_db,
                unit_code="RT101",
                unit_title="Round Trip Unit",
                duration_weeks=1,
            )
            await svc._bg_task

            from app.models.unit import Unit
            from app.services.content_json_renderer import render_content_json

            unit = test_db.query(Unit).filter(Unit.code == "RT101").first()
            assert unit is not None
            mat = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == unit.id)
                .first()
            )
            assert mat is not None and mat.content_json is not None

            markdown = render_content_json(mat.content_json, target="pandoc")
            assert "::: notes" in markdown
            assert "Delivery prompt for the lecturer" in markdown
        finally:
            git_content_service._git_service = original_singleton


# ---------------------------------------------------------------------------
# Phase 2: mixed-format zip imports
# ---------------------------------------------------------------------------


def _html_for_zip() -> bytes:
    return (
        b"<html><body><h1>HTML Lecture</h1>"
        b"<p>HTML body paragraph</p></body></html>"
    )


def _md_for_zip() -> bytes:
    return b"# MD Lecture\n\nMD body paragraph\n"


class TestUnifiedImportPhase2Formats:
    """Bulk LMS imports must route every supported format through the
    structural parser the same way single-file uploads do. Phase 1
    covered PPTX; Phase 2 adds DOCX, HTML, MD, PDF."""

    @pytest.mark.asyncio
    async def test_html_in_zip_produces_structured_content_json(
        self,
        test_db: Session,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        original_singleton = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )
        try:
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w") as zf:
                zf.writestr("lecture_01.html", _html_for_zip())
            buf.seek(0)

            svc = UnifiedImportService()
            svc.apply(
                buf.read(),
                user_id=str(test_user.id),
                user_email=test_user.email,
                db=test_db,
                unit_code="HTML101",
                unit_title="HTML Test",
                duration_weeks=1,
            )
            await svc._bg_task

            from app.models.unit import Unit

            unit = test_db.query(Unit).filter(Unit.code == "HTML101").first()
            assert unit is not None
            mat = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == unit.id)
                .first()
            )
            assert mat is not None
            assert mat.content_json is not None, (
                "HTML in bulk zip should produce structured content_json"
            )
            assert mat.content_json["type"] == "doc"
            assert "HTML body paragraph" in str(mat.content_json)
        finally:
            git_content_service._git_service = original_singleton

    @pytest.mark.asyncio
    async def test_md_in_zip_produces_structured_content_json(
        self,
        test_db: Session,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        original_singleton = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )
        try:
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w") as zf:
                zf.writestr("notes_01.md", _md_for_zip())
            buf.seek(0)

            svc = UnifiedImportService()
            svc.apply(
                buf.read(),
                user_id=str(test_user.id),
                user_email=test_user.email,
                db=test_db,
                unit_code="MD101",
                unit_title="MD Test",
                duration_weeks=1,
            )
            await svc._bg_task

            from app.models.unit import Unit

            unit = test_db.query(Unit).filter(Unit.code == "MD101").first()
            assert unit is not None
            mat = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == unit.id)
                .first()
            )
            assert mat is not None
            assert mat.content_json is not None
            assert mat.content_json["type"] == "doc"
            assert "MD body paragraph" in str(mat.content_json)
        finally:
            git_content_service._git_service = original_singleton

    @pytest.mark.asyncio
    async def test_mixed_format_zip_all_structured(
        self,
        test_db: Session,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        """A zip containing PPTX + HTML + MD should produce three
        structured WeeklyMaterials, all with content_json populated."""
        original_singleton = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )
        try:
            pptx = _build_pptx("PPTX Lecture", "pptx body", "pptx notes")
            buf = io.BytesIO()
            with zipfile.ZipFile(buf, "w") as zf:
                zf.writestr("lecture_01.pptx", pptx)
                zf.writestr("reading_01.html", _html_for_zip())
                zf.writestr("notes_01.md", _md_for_zip())
            buf.seek(0)

            svc = UnifiedImportService()
            svc.apply(
                buf.read(),
                user_id=str(test_user.id),
                user_email=test_user.email,
                db=test_db,
                unit_code="MIX101",
                unit_title="Mixed Unit",
                duration_weeks=1,
            )
            await svc._bg_task

            from app.models.unit import Unit

            unit = test_db.query(Unit).filter(Unit.code == "MIX101").first()
            assert unit is not None
            materials = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == unit.id)
                .all()
            )
            assert len(materials) == 3, (
                f"expected 3 materials, got {len(materials)}"
            )
            # Every material in the mixed zip should have structured
            # content_json — no legacy plain-text fallback for any format
            for mat in materials:
                assert mat.content_json is not None, (
                    f"material '{mat.title}' has no content_json "
                    "(should be structured)"
                )
                assert mat.content_json.get("type") == "doc"
        finally:
            git_content_service._git_service = original_singleton
