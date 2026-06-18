"""
Tests for Mode B (multi-file zip → existing unit) — Phase 3.

Covers:
- ``/batch/preview``: a multi-format zip groups lecture_01.{pptx,pdf,docx}
  into one canonical group with two source files, plus a standalone
  reading.pdf, with unsupported files surfaced as warnings.
- batch apply end-to-end (via the service, awaiting the background task —
  the same pattern test_unified_import_pptx.py uses): asserts the right
  number of materials, content_json populated, source files written to
  git, and ``attached_source_files`` metadata recorded.

Real in-memory SQLite (conftest fixtures); zips built in-memory with
zipfile/io; tiny real PPTX/DOCX fixtures via python-pptx/python-docx.
No mocking except the git singleton is pointed at a tmp dir so writes
don't touch the real content_repos directory.
"""

from __future__ import annotations

import io
import zipfile
from typing import TYPE_CHECKING

import pytest
from docx import Document
from pptx import Presentation

from app.models.weekly_material import WeeklyMaterial
from app.services import git_content_service
from app.services.batch_import_service import BatchImportService
from app.services.material_parsers.persistence import material_source_files_dir

if TYPE_CHECKING:
    from pathlib import Path

    from fastapi.testclient import TestClient
    from sqlalchemy.orm import Session

    from app.models.unit import Unit
    from app.models.user import User


# ---------------------------------------------------------------------------
# Fixture builders
# ---------------------------------------------------------------------------


def _pptx_bytes(title: str = "Lecture One") -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Key point one"
    slide.notes_slide.notes_text_frame.text = "Speak slowly"
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _docx_bytes() -> bytes:
    doc = Document()
    doc.add_heading("Lecture One Handout", level=1)
    doc.add_paragraph("Handout body paragraph.")
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def _pdf_bytes(title: str = "Reading") -> bytes:
    import fitz

    doc = fitz.open()
    doc.set_metadata({"title": title})
    page = doc.new_page()
    page.insert_text((72, 72), f"{title} body paragraph.", fontsize=12)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    buf.seek(0)
    return buf.read()


def _multi_format_zip() -> bytes:
    """lecture_01 in 3 formats + a standalone reading.pdf + an unsupported."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        zf.writestr("lecture_01.pptx", _pptx_bytes("Lecture One"))
        zf.writestr("lecture_01.pdf", _pdf_bytes("Lecture One"))
        zf.writestr("lecture_01.docx", _docx_bytes())
        zf.writestr("reading.pdf", _pdf_bytes("Reading"))
        zf.writestr("notes.txt", b"plain text, unsupported")
        # Hidden/system path that must be skipped silently
        zf.writestr("__MACOSX/lecture_01.pptx", b"junk")
    buf.seek(0)
    return buf.read()


# ---------------------------------------------------------------------------
# /batch/preview
# ---------------------------------------------------------------------------


class TestBatchPreview:
    def test_preview_groups_multi_format(
        self, client: TestClient, test_unit: Unit
    ) -> None:
        files = {"file": ("materials.zip", _multi_format_zip(), "application/zip")}
        response = client.post(
            "/api/import/material/batch/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        assert response.status_code == 200, response.text
        data = response.json()

        # One multi-format group: lecture_01 (pptx canonical, pdf+docx source)
        assert len(data["groups"]) == 1
        group = data["groups"][0]
        assert group["name"] == "lecture_01"
        assert group["canonicalFilename"] == "lecture_01.pptx"
        assert group["canonicalFormat"] == "pptx"
        assert group["parser"] == "pptx_structural"
        assert group["detectedWeek"] == 1  # "lecture_01" → week 1
        src_names = {s["filename"] for s in group["sourceFiles"]}
        assert src_names == {"lecture_01.pdf", "lecture_01.docx"}

        # One standalone: reading.pdf
        assert len(data["standaloneFiles"]) == 1
        assert data["standaloneFiles"][0]["filename"] == "reading.pdf"
        assert data["standaloneFiles"][0]["parser"] == "pdf_paragraphs"

        # Totals
        assert data["totalMaterials"] == 2
        assert data["totalSourceFiles"] == 2

        # Unsupported notes.txt surfaces as a warning; __MACOSX is silent
        warning_text = " ".join(data["warnings"])
        assert "notes.txt" in warning_text
        assert "__MACOSX" not in warning_text

    def test_preview_writes_nothing(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        before = test_db.query(WeeklyMaterial).count()
        files = {"file": ("m.zip", _multi_format_zip(), "application/zip")}
        client.post(
            "/api/import/material/batch/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        assert test_db.query(WeeklyMaterial).count() == before

    def test_preview_rejects_non_zip(self, client: TestClient, test_unit: Unit) -> None:
        files = {"file": ("x.zip", b"not a zip", "application/zip")}
        response = client.post(
            "/api/import/material/batch/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        assert response.status_code == 400
        assert "zip" in response.json()["detail"].lower()

    def test_preview_rejects_unowned_unit(self, client: TestClient) -> None:
        import uuid

        files = {"file": ("m.zip", _multi_format_zip(), "application/zip")}
        response = client.post(
            "/api/import/material/batch/preview",
            files=files,
            data={"unit_id": str(uuid.uuid4())},
        )
        assert response.status_code == 404


# ---------------------------------------------------------------------------
# batch apply (end-to-end via the service, awaiting the background task)
# ---------------------------------------------------------------------------


class TestBatchApply:
    @pytest.mark.asyncio
    async def test_apply_creates_materials_and_source_files(
        self,
        test_db: Session,
        test_unit: Unit,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        original = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )
        try:
            svc = BatchImportService()
            task_id = svc.apply(
                _multi_format_zip(),
                unit=test_unit,
                user_email=test_user.email,
                db=test_db,
            )
            assert task_id
            assert svc._bg_task is not None
            await svc._bg_task

            materials = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == test_unit.id)
                .all()
            )
            # lecture_01 group + reading.pdf standalone = 2 materials
            assert len(materials) == 2

            # The canonical group material is the PPTX
            lecture = next(m for m in materials if "Lecture One" in (m.title or ""))
            assert lecture.content_json is not None
            assert lecture.content_json["type"] == "doc"
            assert "Speak slowly" in str(lecture.content_json)

            # Source files recorded in metadata
            meta = lecture.material_metadata or {}
            attached = meta.get("attached_source_files")
            assert attached is not None
            recorded_formats = {a["format"] for a in attached}
            assert recorded_formats == {"pdf", "docx"}
            for entry in attached:
                assert entry["original_size"] > 0
                assert len(entry["sha256"]) == 64  # full sha256 hex digest

            # Source files actually written to git
            git = git_content_service.get_git_service()
            src_dir = material_source_files_dir(lecture)
            stored = git.list_directory(str(test_unit.id), src_dir)
            assert len(stored) == 2

            # Standalone reading.pdf landed as its own material
            reading = next(m for m in materials if m is not lecture)
            assert reading.content_json is not None
            assert (reading.material_metadata or {}).get(
                "attached_source_files"
            ) is None
        finally:
            git_content_service._git_service = original

    @pytest.mark.asyncio
    async def test_apply_canonical_override(
        self,
        test_db: Session,
        test_unit: Unit,
        test_user: User,
        tmp_path: Path,
    ) -> None:
        original = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )
        try:
            svc = BatchImportService()
            svc.apply(
                _multi_format_zip(),
                unit=test_unit,
                user_email=test_user.email,
                db=test_db,
                overrides={
                    "lecture_01": {
                        "canonical_filename": "lecture_01.docx",
                        "week_number": 5,
                    }
                },
            )
            assert svc._bg_task is not None
            await svc._bg_task

            materials = (
                test_db.query(WeeklyMaterial)
                .filter(WeeklyMaterial.unit_id == test_unit.id)
                .all()
            )
            lecture = next(
                m
                for m in materials
                if (m.material_metadata or {}).get("attached_source_files")
            )
            # Week override applied
            assert lecture.week_number == 5
            # DOCX is now canonical → pptx + pdf are the source files
            attached = lecture.material_metadata["attached_source_files"]
            recorded_formats = {a["format"] for a in attached}
            assert recorded_formats == {"pptx", "pdf"}
        finally:
            git_content_service._git_service = original

    @pytest.mark.asyncio
    async def test_apply_bad_zip_marks_task_failed(
        self,
        test_db: Session,
        test_unit: Unit,
        test_user: User,
    ) -> None:
        from app.services.import_task_store import get_task

        svc = BatchImportService()
        task_id = svc.apply(
            b"not a zip at all",
            unit=test_unit,
            user_email=test_user.email,
            db=test_db,
        )
        assert svc._bg_task is not None
        await svc._bg_task

        task = get_task(task_id)
        assert task is not None
        assert task.status == "failed"
        assert any("Invalid ZIP" in e for e in task.errors)
