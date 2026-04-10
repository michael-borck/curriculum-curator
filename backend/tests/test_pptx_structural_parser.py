"""
Tests for the structural PPTX material parser.

Builds real PPTX files in-memory using python-pptx and asserts the parser
produces the expected content_json structure. Per the project testing
philosophy these are not mocks — the parser runs against actual PPTX bytes
through the actual python-pptx library.
"""

from __future__ import annotations

import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.util import Inches, Pt

from app.services.material_parsers.pptx_structural import PptxStructuralParser

# ---------------------------------------------------------------------------
# Fixture builders
# ---------------------------------------------------------------------------


def _save_to_bytes(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _empty_pptx() -> bytes:
    return _save_to_bytes(Presentation())


def _slide_with_title_and_body(title: str, body: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = title
    body_placeholder = slide.placeholders[1]
    body_placeholder.text = body
    return _save_to_bytes(prs)


def _slide_with_notes(title: str, notes: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.notes_slide.notes_text_frame.text = notes
    return _save_to_bytes(prs)


def _multi_slide_deck(num_slides: int) -> bytes:
    prs = Presentation()
    for i in range(1, num_slides + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i}"
        slide.placeholders[1].text = f"Body of slide {i}"
        slide.notes_slide.notes_text_frame.text = f"Notes for slide {i}"
    return _save_to_bytes(prs)


def _slide_with_table() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only layout
    slide.shapes.title.text = "Table Slide"

    rows, cols = 3, 2
    left = top = Inches(1)
    width = Inches(6)
    height = Inches(2)
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Row 1A"
    table.cell(1, 1).text = "Row 1B"
    table.cell(2, 0).text = "Row 2A"
    table.cell(2, 1).text = "Row 2B"

    return _save_to_bytes(prs)


def _slide_with_picture() -> bytes:
    """Embed a tiny valid PNG (1x1 transparent) into a slide."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Picture Slide"

    # 1x1 transparent PNG
    png_bytes = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000d49444154789c63000100000005000100"
        "0d0a2db40000000049454e44ae426082"
    )
    slide.shapes.add_picture(io.BytesIO(png_bytes), Inches(1), Inches(1))
    return _save_to_bytes(prs)


def _slide_with_bold_italic() -> bytes:
    """Slide with a body paragraph containing bold and italic runs."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Formatting"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = ""  # clear the default empty paragraph
    p = tf.paragraphs[0]

    run1 = p.add_run()
    run1.text = "plain "
    run2 = p.add_run()
    run2.text = "bold "
    run2.font.bold = True
    run3 = p.add_run()
    run3.text = "italic"
    run3.font.italic = True

    return _save_to_bytes(prs)


# ---------------------------------------------------------------------------
# Parser fixture
# ---------------------------------------------------------------------------


@pytest.fixture
def parser() -> PptxStructuralParser:
    return PptxStructuralParser()


# ---------------------------------------------------------------------------
# Helpers for asserting on content_json
# ---------------------------------------------------------------------------


def _node_types(content: list[dict[str, Any]]) -> list[str]:
    return [n.get("type", "") for n in content]


def _find_first(content: list[dict[str, Any]], node_type: str) -> dict[str, Any] | None:
    for n in content:
        if n.get("type") == node_type:
            return n
        children = n.get("content", [])
        if isinstance(children, list):
            found = _find_first(children, node_type)
            if found:
                return found
    return None


def _all_text(node: dict[str, Any]) -> str:
    if node.get("type") == "text":
        return str(node.get("text", ""))
    return "".join(
        _all_text(child) for child in node.get("content", []) if isinstance(child, dict)
    )


# ---------------------------------------------------------------------------
# Title and basic structure
# ---------------------------------------------------------------------------


class TestBasicParsing:
    @pytest.mark.asyncio
    async def test_empty_deck_returns_empty_doc(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_empty_pptx(), "empty.pptx")
        assert result.parser_used == "pptx_structural"
        assert result.content_json["type"] == "doc"
        assert result.content_json["content"] == []
        # Title falls back to filename stem when no slides have content
        assert result.title == "empty"

    @pytest.mark.asyncio
    async def test_title_extracted_as_h1(
        self, parser: PptxStructuralParser
    ) -> None:
        pptx = _slide_with_title_and_body("Introduction", "First point")
        result = await parser.parse(pptx, "deck.pptx")
        # Deck title comes from the first slide's title
        assert result.title == "Introduction"
        # First node is the heading
        first = result.content_json["content"][0]
        assert first["type"] == "heading"
        assert first["attrs"]["level"] == 1
        assert _all_text(first) == "Introduction"

    @pytest.mark.asyncio
    async def test_body_text_extracted_as_paragraph(
        self, parser: PptxStructuralParser
    ) -> None:
        pptx = _slide_with_title_and_body("Title", "Body text here")
        result = await parser.parse(pptx, "deck.pptx")
        types = _node_types(result.content_json["content"])
        assert "heading" in types
        assert "paragraph" in types or "bulletList" in types
        # The body text appears somewhere in the content tree
        all_text = "".join(_all_text(n) for n in result.content_json["content"])
        assert "Body text here" in all_text

    @pytest.mark.asyncio
    async def test_filename_fallback_when_no_title(
        self, parser: PptxStructuralParser
    ) -> None:
        # Slide layout 6 = blank, no title placeholder
        prs = Presentation()
        prs.slides.add_slide(prs.slide_layouts[6])
        result = await parser.parse(_save_to_bytes(prs), "fallback-name.pptx")
        assert result.title == "fallback-name"


# ---------------------------------------------------------------------------
# Multi-slide structure with slide breaks
# ---------------------------------------------------------------------------


class TestSlideBreaks:
    @pytest.mark.asyncio
    async def test_slide_break_between_slides(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_multi_slide_deck(3), "deck.pptx")
        types = _node_types(result.content_json["content"])
        # Three slides, two slide breaks between them
        assert types.count("slideBreak") == 2

    @pytest.mark.asyncio
    async def test_no_leading_slide_break(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_multi_slide_deck(2), "deck.pptx")
        # First node should not be a slideBreak
        assert result.content_json["content"][0]["type"] != "slideBreak"

    @pytest.mark.asyncio
    async def test_no_trailing_slide_break(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_multi_slide_deck(2), "deck.pptx")
        assert result.content_json["content"][-1]["type"] != "slideBreak"


# ---------------------------------------------------------------------------
# Speaker notes — the original motivation for this work
# ---------------------------------------------------------------------------


class TestSpeakerNotes:
    @pytest.mark.asyncio
    async def test_notes_emitted_as_structured_node(
        self, parser: PptxStructuralParser
    ) -> None:
        pptx = _slide_with_notes(
            "Slide A", "Open with the analogy. Pause for questions."
        )
        result = await parser.parse(pptx, "deck.pptx")
        notes_node = _find_first(
            result.content_json["content"], "speakerNotes"
        )
        assert notes_node is not None
        assert "Open with the analogy" in _all_text(notes_node)

    @pytest.mark.asyncio
    async def test_notes_not_flattened_into_body(
        self, parser: PptxStructuralParser
    ) -> None:
        pptx = _slide_with_notes("Slide A", "Hidden delivery prompt")
        result = await parser.parse(pptx, "deck.pptx")
        # The body of the slide (paragraphs/headings) must NOT contain
        # the notes text — it should only appear inside the speakerNotes node
        non_notes_text = ""
        for node in result.content_json["content"]:
            if node.get("type") != "speakerNotes":
                non_notes_text += _all_text(node)
        assert "Hidden delivery prompt" not in non_notes_text

    @pytest.mark.asyncio
    async def test_notes_per_slide(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_multi_slide_deck(3), "deck.pptx")
        # Each slide has its own notes, so we should find three speakerNotes
        # nodes in the document
        all_notes = []

        def _collect(nodes: list[dict[str, Any]]) -> None:
            for n in nodes:
                if n.get("type") == "speakerNotes":
                    all_notes.append(n)
                children = n.get("content", [])
                if isinstance(children, list):
                    _collect(children)

        _collect(result.content_json["content"])
        assert len(all_notes) == 3
        # And each contains its slide-specific text
        all_text = " ".join(_all_text(n) for n in all_notes)
        assert "Notes for slide 1" in all_text
        assert "Notes for slide 2" in all_text
        assert "Notes for slide 3" in all_text

    @pytest.mark.asyncio
    async def test_empty_notes_omitted(
        self, parser: PptxStructuralParser
    ) -> None:
        # Slide with title but no notes content — no speakerNotes node emitted
        pptx = _slide_with_title_and_body("Title", "Body")
        result = await parser.parse(pptx, "deck.pptx")
        notes_node = _find_first(
            result.content_json["content"], "speakerNotes"
        )
        assert notes_node is None


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------


class TestTables:
    @pytest.mark.asyncio
    async def test_table_emitted_as_table_node(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_slide_with_table(), "table.pptx")
        table = _find_first(result.content_json["content"], "table")
        assert table is not None
        # First row should be header cells
        rows = [n for n in table.get("content", []) if n.get("type") == "tableRow"]
        assert len(rows) == 3
        first_row_cells = rows[0]["content"]
        assert all(c["type"] == "tableHeader" for c in first_row_cells)
        # Subsequent rows should be regular cells
        second_row_cells = rows[1]["content"]
        assert all(c["type"] == "tableCell" for c in second_row_cells)
        # Cell text preserved
        assert "Header A" in _all_text(table)
        assert "Row 2B" in _all_text(table)


# ---------------------------------------------------------------------------
# Pictures and image extraction
# ---------------------------------------------------------------------------


class TestImages:
    @pytest.mark.asyncio
    async def test_picture_produces_image_node_and_extracted_image(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_slide_with_picture(), "pic.pptx")
        image_node = _find_first(result.content_json["content"], "image")
        assert image_node is not None
        # The src should be a placeholder filename matching an extracted image
        src = image_node["attrs"]["src"]
        assert src.startswith("slide-1-")
        assert len(result.images) == 1
        assert result.images[0].filename == src
        assert len(result.images[0].data) > 0

    @pytest.mark.asyncio
    async def test_image_filenames_unique_per_slide(
        self, parser: PptxStructuralParser
    ) -> None:
        # Two slides each with one picture — filenames should differ by slide
        prs = Presentation()
        png_bytes = bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
            "890000000d49444154789c63000100000005000100"
            "0d0a2db40000000049454e44ae426082"
        )
        for i in range(1, 3):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Slide {i}"
            slide.shapes.add_picture(
                io.BytesIO(png_bytes), Inches(1), Inches(1)
            )
        result = await parser.parse(_save_to_bytes(prs), "deck.pptx")
        filenames = [img.filename for img in result.images]
        assert len(filenames) == 2
        assert filenames[0] != filenames[1]
        assert "slide-1" in filenames[0]
        assert "slide-2" in filenames[1]


# ---------------------------------------------------------------------------
# Text formatting (bold, italic)
# ---------------------------------------------------------------------------


class TestTextFormatting:
    @pytest.mark.asyncio
    async def test_bold_and_italic_marks(
        self, parser: PptxStructuralParser
    ) -> None:
        result = await parser.parse(_slide_with_bold_italic(), "fmt.pptx")

        # Find any text nodes and collect their marks
        marks_seen: list[str] = []

        def _walk(nodes: list[dict[str, Any]]) -> None:
            for n in nodes:
                if n.get("type") == "text":
                    marks_seen.extend(
                        m.get("type", "") for m in n.get("marks", [])
                    )
                children = n.get("content", [])
                if isinstance(children, list):
                    _walk(children)

        _walk(result.content_json["content"])
        assert "bold" in marks_seen
        assert "italic" in marks_seen
