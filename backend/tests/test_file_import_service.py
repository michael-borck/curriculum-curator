"""
Unit tests for FileImportService (legacy file processing).

Covers text extraction (TXT/MD/DOCX/PDF paths), content-type detection,
section parsing, suggestion/gap generation, filename/week inference,
duration/difficulty heuristics, ZIP analysis, and template stripping.

Test documents are built in-memory with python-docx / python-pptx /
PyPDF2 rather than committed as binary fixtures. PPTX image extraction
is already covered by tests/test_image_upload.py and the structural
parser path by tests/test_unified_import_pptx.py, so neither is
duplicated here.
"""

from __future__ import annotations

import io
import zipfile

import pytest
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfWriter

from app.services.file_import_service import FileImportService

# ---------------------------------------------------------------------------
# Builders
# ---------------------------------------------------------------------------


def _build_docx(
    paragraphs: list[str], table_rows: list[list[str]] | None = None
) -> bytes:
    """Construct a minimal DOCX with paragraphs and an optional table."""
    doc = Document()
    for text in paragraphs:
        doc.add_paragraph(text)
    if table_rows:
        table = doc.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r, row in enumerate(table_rows):
            for c, cell_text in enumerate(row):
                table.rows[r].cells[c].text = cell_text
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _build_pptx(slide_titles: list[str]) -> bytes:
    """Construct a minimal PPTX with one title-only slide per entry."""
    prs = Presentation()
    for title in slide_titles:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
        slide.shapes.title.text = title
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_blank_pdf() -> bytes:
    """Construct a valid single-page PDF with no extractable text."""
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _build_zip(entries: dict[str, bytes]) -> bytes:
    """Build a zip archive from a {path: content} mapping."""
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for path, content in entries.items():
            zf.writestr(path, content)
    return buf.getvalue()


QUIZ_TEXT = (
    "Quiz 1\n"
    "Question 1. What is HTML? (multiple choice)\n"
    "Question 2. CSS cascades. (true/false)\n"
    "Question 3. Select the correct answer below.\n"
)

SLIDES_TEXT = (
    "Lecture 3\n"
    "Learning objectives: understand layout\n"
    "Agenda: today we will cover flexbox\n"
    "In this lecture we explore CSS.\n"
    "Slide 2 introduces the grid presentation.\n"
)

NEUTRAL_TEXT = "The cat sat quietly on the warm mat while dogs ran outside."


@pytest.fixture
def service() -> FileImportService:
    return FileImportService()


# ---------------------------------------------------------------------------
# process_file — extraction paths and result shape
# ---------------------------------------------------------------------------


class TestProcessFile:
    @pytest.mark.asyncio
    async def test_txt_returns_full_result_shape(self, service: FileImportService):
        content = SLIDES_TEXT.encode("utf-8")
        result = await service.process_file(content, "lecture_3.txt", "text/plain")

        assert result["success"] is True
        assert "flexbox" in result["content"]
        assert result["metadata"]["filename"] == "lecture_3.txt"
        assert result["metadata"]["size"] == len(content)
        assert result["metadata"]["extension"] == ".txt"
        assert result["word_count"] == len(SLIDES_TEXT.split())
        assert result["estimated_reading_time"] == result["word_count"] // 200
        assert "slides" in result["available_content_types"]
        assert isinstance(result["sections"], list)
        assert isinstance(result["suggestions"], list)
        assert isinstance(result["gaps"], list)
        assert result["images"] == []

    @pytest.mark.asyncio
    async def test_unsupported_extension_returns_error(
        self, service: FileImportService
    ):
        result = await service.process_file(b"binary blob", "data.xyz", None)

        assert result["success"] is False
        assert "Unsupported file format" in result["error"]
        assert result["content"] == ""
        assert result["content_type"] == "unknown"
        assert result["sections"] == []
        assert result["images"] == []

    @pytest.mark.asyncio
    async def test_empty_txt_file_succeeds_with_general_type(
        self, service: FileImportService
    ):
        result = await service.process_file(b"", "empty.txt", "text/plain")

        assert result["success"] is True
        assert result["content"] == ""
        assert result["content_type"] == "general"
        assert result["content_type_confidence"] == 0.0
        assert result["word_count"] == 0

    @pytest.mark.asyncio
    async def test_docx_extracts_paragraphs_and_tables(
        self, service: FileImportService
    ):
        docx_bytes = _build_docx(
            paragraphs=["Welcome to the unit.", "This covers web development."],
            table_rows=[["Week", "Topic"], ["1", "HTML Basics"]],
        )
        result = await service.process_file(docx_bytes, "notes.docx", None)

        assert result["success"] is True
        assert "Welcome to the unit." in result["content"]
        assert "This covers web development." in result["content"]
        # Table rows are joined with " | "
        assert "Week | Topic" in result["content"]
        assert "1 | HTML Basics" in result["content"]

    @pytest.mark.asyncio
    async def test_markdown_strips_html_tags(self, service: FileImportService):
        md = b"# Course Overview\n\nSome **bold** statement here.\n"
        result = await service.process_file(md, "overview.md", None)

        assert result["success"] is True
        assert "Course Overview" in result["content"]
        assert "bold" in result["content"]
        assert "<" not in result["content"]
        assert "**" not in result["content"]

    @pytest.mark.asyncio
    async def test_blank_pdf_processes_without_error(self, service: FileImportService):
        result = await service.process_file(_build_blank_pdf(), "blank.pdf", None)

        assert result["success"] is True
        assert result["content"] == ""  # no extractable text on a blank page
        assert result["content_type"] == "general"

    @pytest.mark.asyncio
    async def test_corrupt_docx_returns_error_not_exception(
        self, service: FileImportService
    ):
        result = await service.process_file(b"not a real docx", "broken.docx", None)

        assert result["success"] is False
        assert result["error"]

    @pytest.mark.asyncio
    async def test_batch_process_preserves_order(self, service: FileImportService):
        files: list[tuple[bytes, str, str | None]] = [
            (b"first file text", "a.txt", None),
            (b"blob", "b.xyz", None),
        ]
        results = await service.batch_process(files)

        assert len(results) == 2
        assert results[0]["success"] is True
        assert results[0]["metadata"]["filename"] == "a.txt"
        assert results[1]["success"] is False


# ---------------------------------------------------------------------------
# Content-type detection
# ---------------------------------------------------------------------------


class TestDetectContentType:
    def test_quiz_content_detected(self, service: FileImportService):
        detected, confidence, scores = service._detect_content_type(QUIZ_TEXT)

        assert detected == "quiz"
        assert confidence >= 0.3
        assert scores["quiz"] > 0

    def test_slides_content_detected(self, service: FileImportService):
        detected, confidence, _scores = service._detect_content_type(SLIDES_TEXT)

        assert detected == "slides"
        assert confidence >= 0.3

    def test_neutral_content_falls_back_to_general(self, service: FileImportService):
        detected, confidence, _scores = service._detect_content_type(NEUTRAL_TEXT)

        assert detected == "general"
        assert confidence < 0.3


# ---------------------------------------------------------------------------
# Section parsing
# ---------------------------------------------------------------------------


class TestParseSections:
    def test_markdown_headers_split_sections(self, service: FileImportService):
        content = (
            "Some preamble text.\n"
            "# Getting Started\n"
            "First section body.\n"
            "# Advanced Topics\n"
            "Second section body.\n"
        )
        sections = service._parse_sections(content)

        titles = [s["title"] for s in sections]
        assert titles == ["Introduction", "Getting Started", "Advanced Topics"]
        assert sections[0]["content"] == "Some preamble text."
        assert sections[1]["content"] == "First section body."

    def test_chapter_prefix_starts_new_section(self, service: FileImportService):
        content = "Chapter 2: Data Structures\nLists and trees explained here.\n"
        sections = service._parse_sections(content)

        # The chapter prefix and number are preserved in the section title.
        assert sections[0]["title"] == "Chapter 2 Data Structures"
        assert "Lists and trees" in sections[0]["content"]

    def test_empty_content_yields_no_sections(self, service: FileImportService):
        assert service._parse_sections("") == []


# ---------------------------------------------------------------------------
# Suggestions and gap analysis
# ---------------------------------------------------------------------------


class TestSuggestionsAndGaps:
    def test_brief_content_without_objectives_gets_general_suggestions(
        self, service: FileImportService
    ):
        suggestions = service._generate_suggestions(NEUTRAL_TEXT, "general")

        assert any("brief" in s for s in suggestions)
        assert any("learning objectives" in s for s in suggestions)

    def test_quiz_without_answers_suggests_answer_key(self, service: FileImportService):
        suggestions = service._generate_suggestions(
            "Question 1. What is a stack? Question 2. Define a queue.", "quiz"
        )

        assert any("answer key" in s.lower() for s in suggestions)
        assert any("point values" in s.lower() for s in suggestions)

    def test_lab_gap_analysis_flags_missing_safety_as_high(
        self, service: FileImportService
    ):
        gaps = service._analyze_gaps("Mix the solution and record the data.", "lab")

        safety_gaps = [g for g in gaps if g["element"] == "Safety"]
        assert len(safety_gaps) == 1
        assert safety_gaps[0]["severity"] == "high"

    def test_complete_slides_content_has_fewer_gaps(self, service: FileImportService):
        content = (
            "Learning objective: students will be able to build pages. "
            "Introduction and overview of the topic concept. "
            "For example, such as this. Summary and key points. "
            "Activity: discuss and engage. Check your understanding. "
            "Extension challenge for advanced students."
        )
        gaps = service._analyze_gaps(content, "slides")

        assert gaps == []


# ---------------------------------------------------------------------------
# Filename / metadata inference
# ---------------------------------------------------------------------------


class TestNameInference:
    @pytest.mark.parametrize(
        ("filename", "folder", "expected"),
        [
            ("week_3_lecture.pptx", "", 3),
            ("Week 12 Notes.docx", "", 12),
            ("lecture_05.pdf", "", 5),
            ("03_lecture.pdf", "", 3),  # leading number in filename, no folder
            ("notes.txt", "Week_7", 7),
            ("notes.txt", "", None),
            ("week_99_lecture.txt", "", None),  # out of 1-52 range
        ],
    )
    def test_detect_week_number(
        self,
        service: FileImportService,
        filename: str,
        folder: str,
        expected: int | None,
    ):
        assert service.detect_week_number(filename, folder) == expected

    @pytest.mark.parametrize(
        ("filename", "expected"),
        [
            ("Unit_Outline_2026.pdf", "syllabus"),
            ("ISYS1000 Syllabus.docx", "syllabus"),
            ("quiz_week3.docx", "quiz"),
            ("lecture_slides.pptx", "slides"),
            ("worksheet_2.pdf", "worksheet"),
            ("random_notes.txt", "general"),
        ],
    )
    def test_detect_content_type_from_name(
        self, service: FileImportService, filename: str, expected: str
    ):
        assert service.detect_content_type_from_name(filename) == expected

    def test_is_unit_outline(self, service: FileImportService):
        assert service.is_unit_outline("Unit Outline.pdf") is True
        assert service.is_unit_outline("course_guide.docx") is True
        assert service.is_unit_outline("lecture_1.pptx") is False


# ---------------------------------------------------------------------------
# Difficulty, duration, tags, categorisation
# ---------------------------------------------------------------------------


class TestHeuristics:
    @pytest.mark.parametrize(
        ("content", "expected"),
        [
            ("A basic introduction with a simple overview to get started.", "beginner"),
            (
                "Intermediate practice: develop and implement an application.",
                "intermediate",
            ),
            (
                "Advanced research into complex algorithm optimization theory.",
                "advanced",
            ),
        ],
    )
    def test_assess_difficulty(
        self, service: FileImportService, content: str, expected: str
    ):
        assert service.assess_difficulty(content) == expected

    def test_estimate_duration_applies_type_minimums(self, service: FileImportService):
        short = "A few words only."
        assert service.estimate_duration(short, "quiz") == 15
        assert service.estimate_duration(short, "lab") == 45
        assert service.estimate_duration(short, "general") == 10

    def test_estimate_duration_scales_with_word_count(self, service: FileImportService):
        long_content = "word " * 1500
        # 1500 words at 150 wpm for slides = 10 minutes
        assert service.estimate_duration(long_content, "slides") == 10

    def test_extract_tags_filters_stop_words_and_short_words(
        self, service: FileImportService
    ):
        content = "Python python PYTHON is the best. We use python for data and data."
        tags = service._extract_tags(content)

        assert tags[0] == "python"
        assert "data" in tags
        assert "the" not in tags  # stop word
        assert "use" not in tags  # length <= 3

    def test_categorize_content_returns_full_analysis(self, service: FileImportService):
        result = service.categorize_content(QUIZ_TEXT)

        assert result["content_type"] == "quiz"
        assert result["content_type_confidence"] >= 0.3
        assert result["estimated_duration"] >= 15
        assert result["difficulty_level"] in {"beginner", "intermediate", "advanced"}
        assert isinstance(result["tags"], list)
        assert "quiz" not in result["alternative_types"]


# ---------------------------------------------------------------------------
# ZIP processing
# ---------------------------------------------------------------------------


class TestProcessZipFile:
    @pytest.mark.asyncio
    async def test_zip_groups_files_by_week_and_finds_outline(
        self, service: FileImportService
    ):
        zip_bytes = _build_zip(
            {
                "week_1/lecture_notes.txt": SLIDES_TEXT.encode(),
                "week_2/quiz_questions.txt": QUIZ_TEXT.encode(),
                "unit_outline.txt": b"Unit outline for the semester.",
                ".hidden.txt": b"hidden",  # dotfiles must be skipped
                # Files inside hidden/system directories must be skipped too
                "__MACOSX/week_1/resource.txt": b"macos junk",
                ".git/config.txt": b"git junk",
            }
        )
        # db and current_user are accepted but unused by the implementation
        result = await service.process_zip_file(zip_bytes, "unit-1", None, None)

        assert result["unit_outline_found"] is True
        assert result["unit_outline_file"]["filename"] == "unit_outline.txt"
        processed_names = {f["filename"] for f in result["processed_files"]}
        assert ".hidden.txt" not in processed_names
        assert "resource.txt" not in processed_names  # inside __MACOSX/
        assert "config.txt" not in processed_names  # inside .git/
        assert "uploaded.zip" not in processed_names  # the temp zip itself
        assert processed_names == {
            "lecture_notes.txt",
            "quiz_questions.txt",
            "unit_outline.txt",
        }
        assert result["total_files"] == 3
        assert set(result["files_by_week"]) == {1, 2}
        week1 = result["files_by_week"][1][0]
        assert week1["filename"] == "lecture_notes.txt"
        assert week1["detected_type"] == "slides"
        week2 = result["files_by_week"][2][0]
        assert week2["detected_type"] == "quiz"

        # Suggested structure built per week; only slides/worksheet/reading
        # types appear in suggested_content
        structure = {w["week"]: w for w in result["suggested_structure"]}
        assert structure[1]["file_types"] == {"slides": 1}
        assert structure[1]["suggested_content"][0]["type"] == "slides"
        assert structure[2]["file_types"] == {"quiz": 1}
        assert structure[2]["suggested_content"] == []


# ---------------------------------------------------------------------------
# Template stripping (Pandoc reference docs)
# ---------------------------------------------------------------------------


class TestTemplateStripping:
    def test_strip_docx_removes_body_content(self, service: FileImportService):
        docx_bytes = _build_docx(
            paragraphs=["Paragraph one.", "Paragraph two."],
            table_rows=[["a", "b"]],
        )
        stripped = service.strip_docx_to_template(docx_bytes)

        doc = Document(io.BytesIO(stripped))
        assert all(not p.text.strip() for p in doc.paragraphs)
        assert len(doc.tables) == 0

    def test_strip_pptx_removes_all_slides_keeps_masters(
        self, service: FileImportService
    ):
        pptx_bytes = _build_pptx(["Slide One", "Slide Two", "Slide Three"])
        stripped = service.strip_pptx_to_template(pptx_bytes)

        prs = Presentation(io.BytesIO(stripped))
        assert len(prs.slides) == 0
        assert len(prs.slide_masters) >= 1
        assert len(prs.slide_masters[0].slide_layouts) >= 1
