"""
Tests for material image upload, listing, serving, and deletion.

Uses the `client` fixture (TestClient with in-memory DB + auth bypass)
and a temporary content_repos directory.
"""

import io
import uuid

import pytest
from unittest.mock import patch

from app.models.weekly_material import WeeklyMaterial
from app.services.file_import_service import FileImportService
from app.services.git_content_service import GitContentService


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

# Minimal valid 1x1 PNG (67 bytes)
TINY_PNG = (
    b"\x89PNG\r\n\x1a\n"
    b"\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx"
    b"\x9cc\xf8\x0f\x00\x00\x01\x01\x00\x05\x18\xd8N\x00"
    b"\x00\x00\x00IEND\xaeB`\x82"
)


@pytest.fixture
def test_material(test_db, test_unit) -> WeeklyMaterial:
    """Insert a WeeklyMaterial for image tests."""
    mat = WeeklyMaterial(
        id=str(uuid.uuid4()),
        unit_id=test_unit.id,
        week_number=1,
        title="Lecture 1",
        type="lecture",
        description="<p>Hello</p>",
        order_index=0,
    )
    test_db.add(mat)
    test_db.commit()
    test_db.refresh(mat)
    return mat


@pytest.fixture
def git_service(tmp_path):
    """GitContentService backed by a tmp directory."""
    return GitContentService(repos_base=str(tmp_path / "repos"))


# ---------------------------------------------------------------------------
# Unit tests for GitContentService binary methods
# ---------------------------------------------------------------------------


class TestGitContentServiceBinary:
    def test_save_binary_writes_and_commits(self, git_service: GitContentService):
        commit = git_service.save_binary(
            unit_id="unit-1",
            path="images/test.png",
            data=TINY_PNG,
            user_email="test@example.com",
            message="Add test image",
        )
        assert commit  # non-empty commit hash

        # File exists on disk
        data = git_service.read_binary("unit-1", "images/test.png")
        assert data == TINY_PNG

    def test_list_directory(self, git_service: GitContentService):
        git_service.save_binary(
            "unit-1", "imgs/a.png", TINY_PNG, "test@example.com"
        )
        git_service.save_binary(
            "unit-1", "imgs/b.jpg", b"\xff\xd8\xff\xe0", "test@example.com"
        )

        files = git_service.list_directory("unit-1", "imgs")
        assert sorted(files) == ["a.png", "b.jpg"]

    def test_list_directory_nonexistent(self, git_service: GitContentService):
        files = git_service.list_directory("unit-1", "no-such-dir")
        assert files == []

    def test_delete_file(self, git_service: GitContentService):
        git_service.save_binary(
            "unit-1", "imgs/del.png", TINY_PNG, "test@example.com"
        )
        commit = git_service.delete_file(
            "unit-1", "imgs/del.png", "test@example.com"
        )
        assert commit

        with pytest.raises(FileNotFoundError):
            git_service.read_binary("unit-1", "imgs/del.png")

    def test_delete_file_not_found(self, git_service: GitContentService):
        # Ensure repo exists
        git_service.save_binary(
            "unit-1", "imgs/x.png", TINY_PNG, "test@example.com"
        )
        with pytest.raises(FileNotFoundError):
            git_service.delete_file(
                "unit-1", "imgs/nope.png", "test@example.com"
            )

    def test_read_binary_not_found(self, git_service: GitContentService):
        with pytest.raises(FileNotFoundError):
            git_service.read_binary("no-unit", "nope.png")


# ---------------------------------------------------------------------------
# Integration tests via TestClient (HTTP endpoints)
# ---------------------------------------------------------------------------


class TestImageEndpoints:
    """Tests that exercise the full HTTP → route → service chain."""

    def _upload_url(self, unit_id: str, material_id: str) -> str:
        return f"/api/materials/units/{unit_id}/materials/{material_id}/images"

    def _image_url(self, unit_id: str, material_id: str, filename: str) -> str:
        return f"/api/materials/units/{unit_id}/materials/{material_id}/images/{filename}"

    def test_upload_valid_image(self, client, test_unit, test_material, tmp_path):
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=GitContentService(repos_base=str(tmp_path / "repos")),
        ):
            resp = client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("diagram.png", io.BytesIO(TINY_PNG), "image/png")},
            )
        assert resp.status_code == 200
        data = resp.json()
        assert data["filename"] == "diagram.png"
        assert "url" in data

    def test_upload_rejects_non_image(self, client, test_unit, test_material, tmp_path):
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=GitContentService(repos_base=str(tmp_path / "repos")),
        ):
            resp = client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("script.exe", io.BytesIO(b"MZ..."), "application/octet-stream")},
            )
        assert resp.status_code == 400
        assert "Unsupported image type" in resp.json()["detail"]

    def test_upload_rejects_oversized(self, client, test_unit, test_material, tmp_path):
        big_data = b"\x00" * (6 * 1024 * 1024)  # 6MB > 5MB limit
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=GitContentService(repos_base=str(tmp_path / "repos")),
        ):
            resp = client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("big.png", io.BytesIO(big_data), "image/png")},
            )
        assert resp.status_code == 400
        assert "too large" in resp.json()["detail"]

    def test_list_images(self, client, test_unit, test_material, tmp_path):
        svc = GitContentService(repos_base=str(tmp_path / "repos"))
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=svc,
        ):
            # Upload two images
            client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("a.png", io.BytesIO(TINY_PNG), "image/png")},
            )
            client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("b.jpg", io.BytesIO(b"\xff\xd8\xff\xe0"), "image/jpeg")},
            )

            resp = client.get(self._upload_url(test_unit.id, test_material.id))

        assert resp.status_code == 200
        images = resp.json()
        assert len(images) == 2
        filenames = [img["filename"] for img in images]
        assert "a.png" in filenames
        assert "b.jpg" in filenames

    def test_serve_image(self, client, test_unit, test_material, tmp_path):
        svc = GitContentService(repos_base=str(tmp_path / "repos"))
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=svc,
        ):
            client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("photo.png", io.BytesIO(TINY_PNG), "image/png")},
            )

            resp = client.get(
                self._image_url(test_unit.id, test_material.id, "photo.png")
            )

        assert resp.status_code == 200
        assert resp.headers["content-type"] == "image/png"
        assert resp.content == TINY_PNG

    def test_serve_image_not_found(self, client, test_unit, test_material, tmp_path):
        svc = GitContentService(repos_base=str(tmp_path / "repos"))
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=svc,
        ):
            resp = client.get(
                self._image_url(test_unit.id, test_material.id, "nope.png")
            )
        assert resp.status_code == 404

    def test_delete_image(self, client, test_unit, test_material, tmp_path):
        svc = GitContentService(repos_base=str(tmp_path / "repos"))
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=svc,
        ):
            client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={"file": ("del.png", io.BytesIO(TINY_PNG), "image/png")},
            )

            resp = client.delete(
                self._image_url(test_unit.id, test_material.id, "del.png")
            )

        assert resp.status_code == 200
        assert "deleted" in resp.json()["message"].lower()

    def test_delete_image_not_found(self, client, test_unit, test_material, tmp_path):
        svc = GitContentService(repos_base=str(tmp_path / "repos"))
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=svc,
        ):
            resp = client.delete(
                self._image_url(test_unit.id, test_material.id, "nope.png")
            )
        assert resp.status_code == 404

    def test_upload_material_not_found(self, client, test_unit, tmp_path):
        fake_id = str(uuid.uuid4())
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=GitContentService(repos_base=str(tmp_path / "repos")),
        ):
            resp = client.post(
                self._upload_url(test_unit.id, fake_id),
                files={"file": ("x.png", io.BytesIO(TINY_PNG), "image/png")},
            )
        assert resp.status_code == 404

    def test_filename_sanitization(self, client, test_unit, test_material, tmp_path):
        svc = GitContentService(repos_base=str(tmp_path / "repos"))
        with patch(
            "app.api.routes.materials.get_git_service",
            return_value=svc,
        ):
            resp = client.post(
                self._upload_url(test_unit.id, test_material.id),
                files={
                    "file": (
                        "my photo (1).PNG",
                        io.BytesIO(TINY_PNG),
                        "image/png",
                    )
                },
            )
        assert resp.status_code == 200
        fn = resp.json()["filename"]
        # Should be sanitized: no spaces or parens
        assert " " not in fn
        assert "(" not in fn
        assert fn.endswith(".png")


# ---------------------------------------------------------------------------
# PPTX image extraction tests
# ---------------------------------------------------------------------------


def _make_pptx_with_image(image_bytes: bytes, image_ext: str = "png") -> bytes:
    """Create a minimal PPTX file containing one text shape and one picture."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Add a text box
    from pptx.util import Pt

    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    txBox.text_frame.text = "Hello from slide 1"

    # Add a picture from bytes
    img_stream = io.BytesIO(image_bytes)
    slide.shapes.add_picture(img_stream, Inches(1), Inches(2.5), Inches(3), Inches(2))

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class TestPptxImageExtraction:
    """Tests for extracting images from PPTX files."""

    @pytest.mark.asyncio
    async def test_extract_pptx_returns_images(self):
        """_extract_pptx should return image data alongside text."""
        pptx_bytes = _make_pptx_with_image(TINY_PNG, "png")
        service = FileImportService()
        _text, images = await service._extract_pptx(pptx_bytes)

        assert len(images) == 1
        assert images[0]["filename"].startswith("slide-1-")
        assert images[0]["filename"].endswith(".png")
        assert isinstance(images[0]["data"], bytes)
        assert len(images[0]["data"]) > 0

    @pytest.mark.asyncio
    async def test_extract_pptx_placeholder_in_text(self):
        """Extracted text should contain {{IMAGE:filename}} placeholders."""
        pptx_bytes = _make_pptx_with_image(TINY_PNG, "png")
        service = FileImportService()
        text, images = await service._extract_pptx(pptx_bytes)

        filename = images[0]["filename"]
        assert f"{{{{IMAGE:{filename}}}}}" in text

    @pytest.mark.asyncio
    async def test_extract_pptx_text_still_present(self):
        """Text content should still be extracted alongside images."""
        pptx_bytes = _make_pptx_with_image(TINY_PNG, "png")
        service = FileImportService()
        text, _images = await service._extract_pptx(pptx_bytes)

        assert "Hello from slide 1" in text
        assert "--- Slide 1 ---" in text

    @pytest.mark.asyncio
    async def test_process_file_includes_images(self):
        """process_file should include images in the result dict."""
        pptx_bytes = _make_pptx_with_image(TINY_PNG, "png")
        service = FileImportService()
        result = await service.process_file(pptx_bytes, "test.pptx")

        assert result["success"] is True
        assert len(result["images"]) == 1
        assert result["images"][0]["filename"].endswith(".png")

    @pytest.mark.asyncio
    async def test_process_file_non_pptx_has_empty_images(self):
        """Non-PPTX files should have an empty images list."""
        service = FileImportService()
        result = await service.process_file(b"Hello world", "test.txt")

        assert result["success"] is True, f"process_file failed: {result.get('error')}"
        assert result["images"] == []
