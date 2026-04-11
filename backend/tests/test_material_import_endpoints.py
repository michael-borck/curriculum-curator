"""
Tests for the structured material import API (Mode A).

Mode A = single file → existing unit. Phase 1 ships PPTX only via the
``pptx_structural`` parser. The full HTTP → route → parser → service →
DB chain is exercised using the in-memory SQLite client fixture from
conftest.

Image persistence is exercised end-to-end against the real
``GitContentService`` against an in-memory git repo (no mocking).
"""

from __future__ import annotations

import io
import uuid
from pathlib import Path
from typing import TYPE_CHECKING

import pytest
from docx import Document
from pptx import Presentation
from pptx.util import Inches

from app.models.weekly_material import WeeklyMaterial
from app.services.material_parsers.docx_pandoc import DocxPandocParser

if TYPE_CHECKING:
    from fastapi.testclient import TestClient
    from sqlalchemy.orm import Session

    from app.models.unit import Unit


# ---------------------------------------------------------------------------
# PPTX builders (mirror the parser test helpers — kept local so each test
# file is self-contained)
# ---------------------------------------------------------------------------


def _save(prs: Presentation) -> bytes:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def _basic_pptx() -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Imported Lecture"
    slide.placeholders[1].text = "Key point one"
    slide.notes_slide.notes_text_frame.text = "Speak slowly here"
    return _save(prs)


def _multi_slide_with_image_pptx() -> bytes:
    prs = Presentation()
    png = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000d49444154789c63000100000005000100"
        "0d0a2db40000000049454e44ae426082"
    )
    for i in range(1, 3):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = f"Slide {i}"
        slide.shapes.add_picture(io.BytesIO(png), Inches(1), Inches(1))
    return _save(prs)


# ---------------------------------------------------------------------------
# /parsers — registry listing
# ---------------------------------------------------------------------------


class TestParserListing:
    def test_lists_all_phase_2_parsers(self, client: TestClient) -> None:
        response = client.get("/api/import/material/parsers")
        assert response.status_code == 200
        ids = {p["id"] for p in response.json()["parsers"]}
        # Phase 1 + Phase 2 default parsers
        assert ids >= {
            "pptx_structural",
            "docx_pandoc",
            "html_structural",
            "md_structural",
            "pdf_paragraphs",
        }

    def test_pptx_parser_metadata(self, client: TestClient) -> None:
        response = client.get("/api/import/material/parsers")
        # Response is camelCased via CamelModel; field names on the wire
        # use camelCase even though the Python attributes are snake_case
        pptx = next(
            p for p in response.json()["parsers"] if p["id"] == "pptx_structural"
        )
        assert pptx["isDefault"] is True
        assert "pptx" in pptx["supportedFormats"]
        assert pptx["requiresAi"] is False

    def test_phase_2_parsers_are_default_for_their_format(
        self, client: TestClient
    ) -> None:
        response = client.get("/api/import/material/parsers")
        parsers_by_id = {p["id"]: p for p in response.json()["parsers"]}
        expected_defaults = {
            "docx_pandoc": "docx",
            "html_structural": "html",
            "md_structural": "md",
            "pdf_paragraphs": "pdf",
        }
        for parser_id, fmt in expected_defaults.items():
            parser = parsers_by_id[parser_id]
            assert fmt in parser["supportedFormats"]
            assert parser["isDefault"] is True, (
                f"{parser_id} should be the default for .{fmt}"
            )

    def test_filter_by_format(self, client: TestClient) -> None:
        response = client.get("/api/import/material/parsers?format=pptx")
        assert response.status_code == 200
        ids = [p["id"] for p in response.json()["parsers"]]
        assert ids == ["pptx_structural"]

    def test_filter_by_unknown_format(self, client: TestClient) -> None:
        # ipynb has no registered parser yet (Phase 5+)
        response = client.get("/api/import/material/parsers?format=ipynb")
        assert response.status_code == 200
        assert response.json()["parsers"] == []


# ---------------------------------------------------------------------------
# /single/preview — parse without persisting
# ---------------------------------------------------------------------------


class TestPreviewEndpoint:
    def test_preview_returns_structured_content(
        self, client: TestClient, test_unit: Unit
    ) -> None:
        files = {
            "file": ("lecture.pptx", _basic_pptx(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        }
        response = client.post(
            "/api/import/material/single/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        assert response.status_code == 200, response.text
        data = response.json()
        assert data["parserUsed"] == "pptx_structural"
        assert data["title"] == "Imported Lecture"
        assert data["contentJson"]["type"] == "doc"
        # The notes appear in the structured content (not flattened)
        all_text = str(data["contentJson"])
        assert "Speak slowly here" in all_text
        assert "Key point one" in all_text

    def test_preview_does_not_create_material(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        before = test_db.query(WeeklyMaterial).count()
        files = {"file": ("x.pptx", _basic_pptx(), "application/octet-stream")}
        client.post(
            "/api/import/material/single/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        after = test_db.query(WeeklyMaterial).count()
        assert after == before

    def test_preview_rejects_unknown_extension(
        self, client: TestClient, test_unit: Unit
    ) -> None:
        files = {"file": ("notes.txt", b"plain text", "text/plain")}
        response = client.post(
            "/api/import/material/single/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        assert response.status_code == 400
        assert "Unsupported format" in response.json()["detail"]

    def test_preview_rejects_empty_file(
        self, client: TestClient, test_unit: Unit
    ) -> None:
        files = {"file": ("empty.pptx", b"", "application/octet-stream")}
        response = client.post(
            "/api/import/material/single/preview",
            files=files,
            data={"unit_id": str(test_unit.id)},
        )
        assert response.status_code == 400
        assert "empty" in response.json()["detail"].lower()

    def test_preview_rejects_unknown_unit(
        self, client: TestClient
    ) -> None:
        files = {"file": ("x.pptx", _basic_pptx(), "application/octet-stream")}
        response = client.post(
            "/api/import/material/single/preview",
            files=files,
            data={"unit_id": str(uuid.uuid4())},
        )
        assert response.status_code == 404


# ---------------------------------------------------------------------------
# /single/apply — parse and create the material
# ---------------------------------------------------------------------------


class TestApplyEndpoint:
    def test_apply_creates_weekly_material(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        files = {"file": ("lecture.pptx", _basic_pptx(), "application/octet-stream")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 2,
                "type": "lecture",
                "category": "in_class",
            },
        )
        assert response.status_code == 200, response.text
        data = response.json()
        assert data["parserUsed"] == "pptx_structural"
        assert data["weekNumber"] == 2
        assert data["title"] == "Imported Lecture"

        # Material exists in DB with content_json populated
        mat = (
            test_db.query(WeeklyMaterial)
            .filter(WeeklyMaterial.id == data["materialId"])
            .first()
        )
        assert mat is not None
        assert mat.content_json is not None
        assert mat.content_json["type"] == "doc"
        # Speaker notes survived as structured content (not flattened)
        all_text = str(mat.content_json)
        assert "Speak slowly here" in all_text

    def test_apply_with_title_override(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        files = {"file": ("x.pptx", _basic_pptx(), "application/octet-stream")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "title_override": "Custom Title",
                "type": "lecture",
            },
        )
        assert response.status_code == 200
        assert response.json()["title"] == "Custom Title"

    def test_apply_persists_images_and_rewrites_src(
        self,
        client: TestClient,
        test_unit: Unit,
        test_db: Session,
        tmp_path: Path,
    ) -> None:
        # Point the git content service at a temp dir for this test so
        # writes don't pollute the real content_repos directory. We patch
        # the module-level singleton to a fresh instance pointed at the
        # temp path, then restore afterwards.
        from app.services import git_content_service

        original_singleton = git_content_service._git_service
        git_content_service._git_service = git_content_service.GitContentService(
            repos_base=str(tmp_path)
        )

        try:
            files = {
                "file": (
                    "deck.pptx",
                    _multi_slide_with_image_pptx(),
                    "application/octet-stream",
                )
            }
            response = client.post(
                "/api/import/material/single/apply",
                files=files,
                data={
                    "unit_id": str(test_unit.id),
                    "week_number": 1,
                    "type": "lecture",
                },
            )
        finally:
            git_content_service._git_service = original_singleton

        assert response.status_code == 200, response.text
        data = response.json()
        assert data["imageCount"] == 2

        mat = (
            test_db.query(WeeklyMaterial)
            .filter(WeeklyMaterial.id == data["materialId"])
            .first()
        )
        assert mat is not None
        # Image src in content_json should now be canonical URLs, not the
        # parser's bare-filename placeholders
        content_str = str(mat.content_json)
        assert "/api/materials/units/" in content_str
        assert "/images/" in content_str
        # Bare placeholders should be gone — no leftover "src': 'slide-N-..."
        assert "'src': 'slide-" not in content_str


# ---------------------------------------------------------------------------
# Round-trip: PPTX → import → export PPTX → notes preserved
# ---------------------------------------------------------------------------


class TestNotesRoundTrip:
    """End-to-end verification that the work in this PR closes the original
    speaker notes round-trip gap. Imports a PPTX with known notes, then runs
    the export pipeline against the resulting WeeklyMaterial and asserts the
    pandoc-bound markdown contains ``::: notes`` fenced divs with the right
    text. We don't actually invoke Pandoc here — that's covered by the
    speaker-notes export tests — but we verify the import-side produces
    structured notes that the export side will route correctly."""

    def test_imported_notes_reach_pandoc_target(
        self,
        client: TestClient,
        test_unit: Unit,
        test_db: Session,
    ) -> None:
        # Step 1: import the PPTX through the new endpoint
        pptx = _basic_pptx()  # has "Speak slowly here" in notes
        files = {"file": ("lecture.pptx", pptx, "application/octet-stream")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "type": "lecture",
            },
        )
        assert response.status_code == 200
        material_id = response.json()["materialId"]

        # Step 2: feed the imported material through the export pipeline
        # using the same target="pandoc" path that PPTX export uses
        from app.services.content_json_renderer import render_content_json

        mat = (
            test_db.query(WeeklyMaterial)
            .filter(WeeklyMaterial.id == material_id)
            .first()
        )
        assert mat is not None
        assert mat.content_json is not None

        markdown = render_content_json(mat.content_json, target="pandoc")

        # The notes survived as a structured node and now render as a
        # Pandoc fenced div that PowerPoint export will route to the
        # speaker notes pane
        assert "::: notes" in markdown
        assert "Speak slowly here" in markdown
        assert "\n\n::: notes\n" in markdown  # blank-line padding for Pandoc
        assert "\n:::\n\n" in markdown


# ---------------------------------------------------------------------------
# Phase 2: DOCX, HTML, MD, PDF via the apply endpoint
# ---------------------------------------------------------------------------


def _docx_bytes() -> bytes:
    doc = Document()
    doc.add_heading("Imported DOCX", level=1)
    doc.add_paragraph("Body paragraph from a Word document.")
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.read()


def _html_bytes() -> bytes:
    return (
        b"<html><body><h1>Imported HTML</h1>"
        b"<p>Body paragraph from an HTML file.</p></body></html>"
    )


def _md_bytes() -> bytes:
    return b"# Imported Markdown\n\nBody paragraph from a markdown file.\n"


def _pdf_bytes() -> bytes:
    import fitz

    doc = fitz.open()
    doc.set_metadata({"title": "Imported PDF"})
    page = doc.new_page()
    page.insert_text((72, 72), "Body paragraph from a PDF.", fontsize=12)
    buf = io.BytesIO()
    doc.save(buf)
    doc.close()
    buf.seek(0)
    return buf.read()


class TestPhase2FormatsApply:
    """Each Phase 2 format should land in the DB as a structured
    WeeklyMaterial via the Mode A apply endpoint, same as PPTX does."""

    @pytest.mark.skipif(
        not DocxPandocParser.is_available(),
        reason="pandoc binary not available",
    )
    def test_docx_apply_creates_structured_material(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        files = {"file": ("lecture.docx", _docx_bytes(), "application/octet-stream")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "type": "lecture",
            },
        )
        assert response.status_code == 200, response.text
        data = response.json()
        assert data["parserUsed"] == "docx_pandoc"
        assert data["title"] == "Imported DOCX"

        mat = (
            test_db.query(WeeklyMaterial)
            .filter(WeeklyMaterial.id == data["materialId"])
            .first()
        )
        assert mat is not None
        assert mat.content_json is not None
        assert "Body paragraph" in str(mat.content_json)

    def test_html_apply_creates_structured_material(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        files = {"file": ("page.html", _html_bytes(), "text/html")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "type": "reading",
            },
        )
        assert response.status_code == 200, response.text
        data = response.json()
        assert data["parserUsed"] == "html_structural"
        assert data["title"] == "Imported HTML"

        mat = (
            test_db.query(WeeklyMaterial)
            .filter(WeeklyMaterial.id == data["materialId"])
            .first()
        )
        assert mat is not None
        assert mat.content_json is not None
        assert "Body paragraph" in str(mat.content_json)

    def test_md_apply_creates_structured_material(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        files = {"file": ("notes.md", _md_bytes(), "text/markdown")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "type": "reading",
            },
        )
        assert response.status_code == 200, response.text
        data = response.json()
        assert data["parserUsed"] == "md_structural"
        assert data["title"] == "Imported Markdown"

    def test_pdf_apply_creates_structured_material(
        self, client: TestClient, test_unit: Unit, test_db: Session
    ) -> None:
        files = {"file": ("paper.pdf", _pdf_bytes(), "application/pdf")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "type": "reading",
            },
        )
        assert response.status_code == 200, response.text
        data = response.json()
        assert data["parserUsed"] == "pdf_paragraphs"
        assert data["title"] == "Imported PDF"

        # PDF always surfaces the "structure not recovered" warning
        assert any(
            "plain paragraphs" in w.lower() for w in data["warnings"]
        )

    def test_unsupported_format_rejected(
        self, client: TestClient, test_unit: Unit
    ) -> None:
        files = {"file": ("notes.txt", b"plain text content", "text/plain")}
        response = client.post(
            "/api/import/material/single/apply",
            files=files,
            data={
                "unit_id": str(test_unit.id),
                "week_number": 1,
                "type": "reading",
            },
        )
        assert response.status_code == 400
        assert ".txt" in response.json()["detail"].lower()
